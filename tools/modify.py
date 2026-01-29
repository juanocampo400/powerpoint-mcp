"""
Modification tools: modify shapes, delete shapes, find and replace text.
"""
from copy import deepcopy
from lxml import etree
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# XML namespace for DrawingML
DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# Bullet type mappings
BULLET_CHARS = {
    "bullet": "\u2022",      # •
    "dash": "\u2013",        # –
    "arrow": "\u2192",       # →
    "check": "\u2713",       # ✓
    "square": "\u25A0",      # ■
    "circle": "\u25CF",      # ●
    "diamond": "\u25C6",     # ◆
    "star": "\u2605",        # ★
}

NUMBERED_TYPES = {
    "number": "arabicPeriod",       # 1. 2. 3.
    "number_paren": "arabicParenR", # 1) 2) 3)
    "roman": "romanLcPeriod",       # i. ii. iii.
    "roman_upper": "romanUcPeriod", # I. II. III.
    "letter": "alphaLcPeriod",      # a. b. c.
    "letter_upper": "alphaUcPeriod", # A. B. C.
}

ALL_BULLET_TYPES = list(BULLET_CHARS.keys()) + list(NUMBERED_TYPES.keys()) + ["none"]


def _process_text_escapes(text: str) -> str:
    """Convert escape sequences like \\n to actual newlines.

    Args:
        text: Input text that may contain escape sequences

    Returns:
        Text with escape sequences converted to actual characters
    """
    if text is None:
        return None
    return text.replace("\\n", "\n").replace("\\t", "\t")


def _capture_paragraph_format(paragraph):
    """Capture formatting properties from a paragraph.

    Args:
        paragraph: A python-pptx Paragraph object

    Returns:
        Dictionary containing captured formatting properties including XML for bullets
    """
    format_dict = {
        'level': paragraph.level,
        'alignment': paragraph.alignment,
        'pPr_xml': None,
    }

    # Capture the paragraph properties XML (includes bullet formatting)
    p_elem = paragraph._p
    pPr = p_elem.find(qn('a:pPr'))
    if pPr is not None:
        # Deep copy the entire pPr element to preserve all formatting
        format_dict['pPr_xml'] = deepcopy(pPr)

    return format_dict


def _apply_paragraph_format(paragraph, format_dict):
    """Apply captured formatting to a paragraph.

    Args:
        paragraph: A python-pptx Paragraph object
        format_dict: Dictionary from _capture_paragraph_format
    """
    # Apply level
    if format_dict.get('level') is not None:
        paragraph.level = format_dict['level']

    # Apply alignment
    if format_dict.get('alignment') is not None:
        paragraph.alignment = format_dict['alignment']

    # Apply paragraph properties XML (includes bullets)
    if format_dict.get('pPr_xml') is not None:
        p_elem = paragraph._p

        # Remove existing pPr if present
        existing_pPr = p_elem.find(qn('a:pPr'))
        if existing_pPr is not None:
            p_elem.remove(existing_pPr)

        # Insert the copied pPr at the beginning
        p_elem.insert(0, deepcopy(format_dict['pPr_xml']))


def _apply_bullet_style(paragraph, bullet_type):
    """Apply a specific bullet style to a paragraph.

    Args:
        paragraph: A python-pptx Paragraph object
        bullet_type: One of BULLET_CHARS keys, NUMBERED_TYPES keys, or "none"
    """
    p_elem = paragraph._p

    # Get or create pPr element
    pPr = p_elem.find(qn('a:pPr'))
    if pPr is None:
        pPr = etree.Element(qn('a:pPr'))
        p_elem.insert(0, pPr)

    # Remove existing bullet elements
    for elem_name in ['a:buNone', 'a:buChar', 'a:buAutoNum']:
        existing = pPr.find(qn(elem_name))
        if existing is not None:
            pPr.remove(existing)

    if bullet_type == "none":
        # Explicitly no bullets
        buNone = etree.SubElement(pPr, qn('a:buNone'))
    elif bullet_type in BULLET_CHARS:
        # Character bullet
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', BULLET_CHARS[bullet_type])
    elif bullet_type in NUMBERED_TYPES:
        # Numbered list
        buAutoNum = etree.SubElement(pPr, qn('a:buAutoNum'))
        buAutoNum.set('type', NUMBERED_TYPES[bullet_type])


def _update_text_preserve_formatting(text_frame, new_text):
    """Update text frame content while preserving paragraph formatting.

    Args:
        text_frame: A python-pptx TextFrame object
        new_text: New text content (newlines separate paragraphs)
    """
    lines = new_text.split('\n') if new_text else ['']
    existing_paras = list(text_frame.paragraphs)

    # Capture formatting from existing paragraphs
    formats = []
    for para in existing_paras:
        formats.append(_capture_paragraph_format(para))

    # Get the last format to use for any extra lines
    last_format = formats[-1] if formats else None

    # Clear all existing paragraphs except the first
    # (TextFrame always has at least one paragraph)
    p_elements = text_frame._txBody.findall(qn('a:p'))
    for p_elem in p_elements[1:]:
        text_frame._txBody.remove(p_elem)

    # Update paragraphs
    for i, line in enumerate(lines):
        if i == 0:
            # Use the first (always existing) paragraph
            para = text_frame.paragraphs[0]
            # Clear existing runs
            _clear_paragraph_runs(para)
        else:
            # Create new paragraph using python-pptx API
            para = text_frame.add_paragraph()

        # Add the text as a run
        run = para.add_run()
        run.text = line

        # Apply formatting from original paragraph (or last available)
        format_to_apply = formats[i] if i < len(formats) else last_format
        if format_to_apply:
            _apply_paragraph_format(para, format_to_apply)


def _update_text_with_bullets(text_frame, new_text, bullet_type):
    """Update text frame with explicit bullet formatting while preserving other formatting.

    Args:
        text_frame: A python-pptx TextFrame object
        new_text: New text content (newlines separate paragraphs)
        bullet_type: Bullet type to apply
    """
    lines = new_text.split('\n') if new_text else ['']
    existing_paras = list(text_frame.paragraphs)

    # Capture formatting from existing paragraphs
    formats = []
    for para in existing_paras:
        formats.append(_capture_paragraph_format(para))

    last_format = formats[-1] if formats else None

    # Clear all existing paragraphs except the first
    p_elements = text_frame._txBody.findall(qn('a:p'))
    for p_elem in p_elements[1:]:
        text_frame._txBody.remove(p_elem)

    # Update paragraphs
    for i, line in enumerate(lines):
        if i == 0:
            para = text_frame.paragraphs[0]
            _clear_paragraph_runs(para)
        else:
            # Create new paragraph using python-pptx API
            para = text_frame.add_paragraph()

        # Add the text as a run
        run = para.add_run()
        run.text = line

        # Apply base formatting from template
        format_to_apply = formats[i] if i < len(formats) else last_format
        if format_to_apply:
            _apply_paragraph_format(para, format_to_apply)

        # Override with explicit bullet style
        _apply_bullet_style(para, bullet_type)


def _clear_paragraph_runs(paragraph):
    """Remove all runs from a paragraph.

    Args:
        paragraph: A python-pptx Paragraph object
    """
    p_elem = paragraph._p
    for r_elem in p_elem.findall(qn('a:r')):
        p_elem.remove(r_elem)
    # Also remove any direct text elements
    for t_elem in p_elem.findall(qn('a:t')):
        p_elem.remove(t_elem)


def _find_table_shape(slide, slide_number, table_index=1, shape_id=None, shape_name=None):
    """Find a table shape on a slide by ID, name, or index.

    Args:
        slide: The slide object
        slide_number: Slide number (for error messages)
        table_index: 1-based index of table (default 1)
        shape_id: Shape ID (takes precedence)
        shape_name: Shape name (takes precedence over table_index)

    Returns:
        Tuple of (table_shape, error_message). One will be None.
    """
    if shape_id is not None:
        for s in slide.shapes:
            if s.shape_id == shape_id:
                if s.has_table:
                    return s, None
                return None, f"Error: Shape with ID {shape_id} is not a table"
        return None, f"Error: Shape with ID {shape_id} not found on slide {slide_number}"

    if shape_name is not None:
        for s in slide.shapes:
            if s.name == shape_name:
                if s.has_table:
                    return s, None
                return None, f"Error: Shape named '{shape_name}' is not a table"
        return None, f"Error: Shape named '{shape_name}' not found on slide {slide_number}"

    # Find by table_index - sort by position for intuitive ordering
    tables = sorted(
        [s for s in slide.shapes if s.has_table],
        key=lambda s: (s.top, s.left)
    )
    if not tables:
        return None, f"Error: No tables found on slide {slide_number}"
    if table_index < 1 or table_index > len(tables):
        return None, f"Error: table_index {table_index} is out of range (1-{len(tables)})"

    return tables[table_index - 1], None


def register_modify_tools(mcp, state):
    """Register modification tools with the MCP server."""

    @mcp.tool()
    def modify_shape(
        slide_number: int,
        shape_id: int = None,
        shape_name: str = None,
        left: float = None,
        top: float = None,
        width: float = None,
        height: float = None,
        text: str = None,
        fill_color: str = None,
        line_color: str = None,
        rotation: float = None,
        bullets: str = None
    ) -> str:
        """
        Modify an existing shape's properties.

        FORMATTING Note: When setting the text parameter,
        paragraph-level formatting (bullets, alignment, indentation) is preserved,
        but run-level formatting (font name, size, color, bold, italic) is reset to defaults.
        for updating styled text content, prefer find_and_repllace instead.

        Args:
            slide_number: Target slide number (1-based)
            shape_id: ID of the shape to modify (use get_slide_snapshot to find IDs)
            shape_name: Name of the shape to modify (alternative to shape_id)
            left: New left position in inches (optional)
            top: New top position in inches (optional)
            width: New width in inches (optional)
            height: New height in inches (optional)
            text: New text content (only for shapes with text frames)
            fill_color: New fill color as hex code
            line_color: New line/border color as hex code
            rotation: Rotation angle in degrees
            bullets: Apply bullet/list formatting to all paragraphs. Options:
                - Character bullets: "bullet", "dash", "arrow", "check", "square", "circle", "diamond", "star"
                - Numbered lists: "number" (1. 2. 3.), "number_paren" (1) 2) 3)),
                  "roman" (i. ii. iii.), "roman_upper" (I. II. III.),
                  "letter" (a. b. c.), "letter_upper" (A. B. C.)
                - "none" to remove existing bullets
                - None (default) to keep existing formatting

        Note: Provide either shape_id OR shape_name. If both provided, shape_id takes precedence.

        Returns:
            Success message with updated properties, or error message
        """
        if state.presentation is None:
            return "Error: No presentation is currently open"

        if shape_id is None and shape_name is None:
            return "Error: Must provide either shape_id or shape_name"

        prs = state.presentation
        if slide_number < 1 or slide_number > len(prs.slides):
            return f"Error: slide_number {slide_number} is out of range (1-{len(prs.slides)})"

        slide = prs.slides[slide_number - 1]

        # Find the shape
        shape = None
        if shape_id:
            for s in slide.shapes:
                if s.shape_id == shape_id:
                    shape = s
                    break
            if not shape:
                return f"Error: Shape with ID {shape_id} not found on slide {slide_number}"
        else:
            for s in slide.shapes:
                if s.name == shape_name:
                    shape = s
                    break
            if not shape:
                return f"Error: Shape named '{shape_name}' not found on slide {slide_number}"

        changes = []

        try:
            # Position changes
            if left is not None:
                shape.left = Inches(left)
                changes.append(f"left={left}\"")
            if top is not None:
                shape.top = Inches(top)
                changes.append(f"top={top}\"")
            if width is not None:
                shape.width = Inches(width)
                changes.append(f"width={width}\"")
            if height is not None:
                shape.height = Inches(height)
                changes.append(f"height={height}\"")

            # Text
            if text is not None:
                if shape.has_text_frame:
                    processed_text = _process_text_escapes(text)

                    # Validate bullet type if specified
                    if bullets is not None and bullets not in ALL_BULLET_TYPES:
                        return f"Error: Invalid bullet type '{bullets}'. Valid types: {', '.join(ALL_BULLET_TYPES)}"

                    if bullets is None:
                        # Preserve template formatting
                        _update_text_preserve_formatting(shape.text_frame, processed_text)
                        changes.append("text updated (formatting preserved)")
                    else:
                        # Apply explicit bullet style (still preserves other formatting)
                        _update_text_with_bullets(shape.text_frame, processed_text, bullets)
                        changes.append(f"text updated (bullets: {bullets})")
                else:
                    return f"Error: Shape {shape.shape_id} does not support text"

            # Bullets only (no text change)
            elif bullets is not None:
                if bullets not in ALL_BULLET_TYPES:
                    return f"Error: Invalid bullet type '{bullets}'. Valid types: {', '.join(ALL_BULLET_TYPES)}"
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        _apply_bullet_style(para, bullets)
                    changes.append(f"bullets: {bullets}")
                else:
                    return f"Error: Shape {shape.shape_id} does not support text"

            # Fill color
            if fill_color is not None:
                color = _parse_color(fill_color)
                if color:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = color
                    changes.append(f"fill={fill_color}")

            # Line color
            if line_color is not None:
                color = _parse_color(line_color)
                if color:
                    shape.line.color.rgb = color
                    changes.append(f"line={line_color}")

            # Rotation
            if rotation is not None:
                shape.rotation = rotation
                changes.append(f"rotation={rotation}°")

            if changes:
                state.is_modified = True
                return f"Successfully modified shape {shape.shape_id} on slide {slide_number}\nChanges: {', '.join(changes)}"
            else:
                return "No changes specified"

        except Exception as e:
            return f"Error modifying shape: {str(e)}"

    @mcp.tool()
    def delete_shape(
        slide_number: int,
        shape_id: int = None,
        shape_name: str = None
    ) -> str:
        """
        Delete a shape from a slide.

        Args:
            slide_number: Target slide number (1-based)
            shape_id: ID of the shape to delete (use get_slide_snapshot to find IDs)
            shape_name: Name of the shape to delete (alternative to shape_id)

        Note: Provide either shape_id OR shape_name. If both provided, shape_id takes precedence.

        Returns:
            Success message, or error message
        """
        if state.presentation is None:
            return "Error: No presentation is currently open"

        if shape_id is None and shape_name is None:
            return "Error: Must provide either shape_id or shape_name"

        prs = state.presentation
        if slide_number < 1 or slide_number > len(prs.slides):
            return f"Error: slide_number {slide_number} is out of range (1-{len(prs.slides)})"

        slide = prs.slides[slide_number - 1]

        # Find the shape
        shape = None
        if shape_id:
            for s in slide.shapes:
                if s.shape_id == shape_id:
                    shape = s
                    break
            if not shape:
                return f"Error: Shape with ID {shape_id} not found on slide {slide_number}"
        else:
            for s in slide.shapes:
                if s.name == shape_name:
                    shape = s
                    break
            if not shape:
                return f"Error: Shape named '{shape_name}' not found on slide {slide_number}"

        try:
            # Get shape element and remove it
            sp = shape._element
            sp.getparent().remove(sp)

            state.is_modified = True
            return f"Successfully deleted shape (ID: {shape_id or 'N/A'}, Name: '{shape_name or 'N/A'}') from slide {slide_number}"
        except Exception as e:
            return f"Error deleting shape: {str(e)}"

    @mcp.tool()
    def find_and_replace(
        find_text: str,
        replace_text: str,
        slide_number: int = None,
        match_case: bool = False
    ) -> str:
        """
        Find and replace text across the presentation or a specific slide.

        PRESERVES FORMATTING: operates at run level, keeps font properties (name,
        size, color, bold, italic) intact for both text frames and table cells.
        Prefer this over modify_shape when updating text in styled templates.

        Args:
            find_text: Text to search for
            replace_text: Text to replace with
            slide_number: Optional - limit to specific slide (1-based). If not provided, searches all slides.
            match_case: Whether to match case exactly (default False)

        Returns:
            Summary of replacements made
        """
        if state.presentation is None:
            return "Error: No presentation is currently open"

        prs = state.presentation

        # Determine which slides to search
        if slide_number is not None:
            if slide_number < 1 or slide_number > len(prs.slides):
                return f"Error: slide_number {slide_number} is out of range (1-{len(prs.slides)})"
            slides_to_search = [(slide_number, prs.slides[slide_number - 1])]
        else:
            slides_to_search = [(i+1, slide) for i, slide in enumerate(prs.slides)]

        replacements = []

        for slide_num, slide in slides_to_search:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            original_text = run.text
                            if match_case:
                                if find_text in original_text:
                                    run.text = original_text.replace(find_text, replace_text)
                                    replacements.append(f"Slide {slide_num}, Shape '{shape.name}'")
                            else:
                                if find_text.lower() in original_text.lower():
                                    # Case-insensitive replace
                                    import re
                                    run.text = re.sub(
                                        re.escape(find_text),
                                        replace_text,
                                        original_text,
                                        flags=re.IGNORECASE
                                    )
                                    replacements.append(f"Slide {slide_num}, Shape '{shape.name}'")

                # Also check tables
                if shape.has_table:
                    for row_idx, row in enumerate(shape.table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    original_text = run.text
                                    if match_case:
                                        if find_text in original_text:
                                            run.text = original_text.replace(find_text, replace_text)
                                            replacements.append(f"Slide {slide_num}, Table row {row_idx+1} col {col_idx+1}")
                                    else:
                                        if find_text.lower() in original_text.lower():
                                            import re
                                            run.text = re.sub(
                                                re.escape(find_text),
                                                replace_text,
                                                original_text,
                                                flags=re.IGNORECASE
                                            )
                                            replacements.append(f"Slide {slide_num}, Table row {row_idx+1} col {col_idx+1}")

        if replacements:
            state.is_modified = True
            return f"Replaced '{find_text}' with '{replace_text}' in {len(replacements)} location(s):\n" + "\n".join(f"  - {r}" for r in replacements[:20]) + ("\n  ..." if len(replacements) > 20 else "")
        else:
            return f"No occurrences of '{find_text}' found"

    @mcp.tool()
    def get_table_content(
        slide_number: int,
        table_index: int = 1,
        shape_id: int = None,
        shape_name: str = None
    ) -> str:
        """
        Get the content of a specific table.

        Args:
            slide_number: Target slide number (1-based)
            table_index: 1-based index of the table on the slide (default 1, first table).
                         Tables are ordered top-to-bottom, left-to-right by position.
            shape_id: ID of the table shape (alternative to table_index)
            shape_name: Name of the table shape (alternative to table_index)

        Note: If shape_id or shape_name is provided, table_index is ignored.
              shape_id takes precedence over shape_name.

        Returns:
            Table content as formatted text with row/column structure.
            Newlines within cells are shown as \\n.
        """
        if state.presentation is None:
            return "Error: No presentation is currently open"

        prs = state.presentation
        if slide_number < 1 or slide_number > len(prs.slides):
            return f"Error: slide_number {slide_number} is out of range (1-{len(prs.slides)})"

        slide = prs.slides[slide_number - 1]

        table_shape, error = _find_table_shape(slide, slide_number, table_index, shape_id, shape_name)
        if error:
            return error

        table = table_shape.table
        rows = len(table.rows)
        cols = len(table.columns)

        # Build output
        result = [f"Table '{table_shape.name}' (ID: {table_shape.shape_id})"]
        result.append(f"Dimensions: {rows} rows x {cols} columns\n")

        for row_idx, row in enumerate(table.rows):
            row_cells = []
            for col_idx, cell in enumerate(row.cells):
                cell_text = cell.text.replace('\n', '\\n')  # Escape newlines for display
                row_cells.append(cell_text)
            result.append(f"Row {row_idx + 1}: {row_cells}")

        return "\n".join(result)

    @mcp.tool()
    def modify_table_cell(
        slide_number: int,
        row: int,
        column: int,
        text: str,
        table_index: int = 1,
        shape_id: int = None,
        shape_name: str = None
    ) -> str:
        """
        Modify the content of a specific table cell with formatting preservation.

        PRESERVES FORMATTING: If the cell has existing content with runs, the text
        of the first run is replaced, preserving font properties (name, size, color,
        bold, italic). For cells without runs, a new run is added to preserve
        paragraph-level formatting.

        Note on multi-paragraph cells: Only the first paragraph is kept; additional
        paragraphs are removed. Use find_and_replace for finer control.

        Note on merged cells: If the cell at (row, column) is part of a merged range,
        python-pptx will modify the merge origin cell. This may affect a larger area
        than expected.

        Args:
            slide_number: Target slide number (1-based)
            row: 1-based row number
            column: 1-based column number
            text: New cell content (use \\n for newlines within the cell)
            table_index: 1-based index of the table on the slide (default 1).
                         Tables are ordered top-to-bottom, left-to-right by position.
            shape_id: ID of the table shape (alternative to table_index)
            shape_name: Name of the table shape (alternative to table_index)

        Note: If shape_id or shape_name is provided, table_index is ignored.
              shape_id takes precedence over shape_name.

        Returns:
            Success message with modification details, or error message
        """
        if state.presentation is None:
            return "Error: No presentation is currently open"

        prs = state.presentation
        if slide_number < 1 or slide_number > len(prs.slides):
            return f"Error: slide_number {slide_number} is out of range (1-{len(prs.slides)})"

        slide = prs.slides[slide_number - 1]

        table_shape, error = _find_table_shape(slide, slide_number, table_index, shape_id, shape_name)
        if error:
            return error

        table = table_shape.table
        num_rows = len(table.rows)
        num_cols = len(table.columns)

        # Validate row/column
        if row < 1 or row > num_rows:
            return f"Error: row {row} is out of range (1-{num_rows})"
        if column < 1 or column > num_cols:
            return f"Error: column {column} is out of range (1-{num_cols})"

        cell = table.cell(row - 1, column - 1)  # Convert to 0-based
        processed_text = _process_text_escapes(text)

        try:
            # Preserve formatting by modifying at run level
            text_frame = cell.text_frame
            paragraphs = list(text_frame.paragraphs)

            if paragraphs:
                first_para = paragraphs[0]
                runs = list(first_para.runs)

                if runs:
                    # Cell has existing runs - modify first run to preserve its formatting
                    runs[0].text = processed_text

                    # Clear any additional runs in first paragraph (list already materialized)
                    for run in runs[1:]:
                        run._r.getparent().remove(run._r)
                else:
                    # Paragraph exists but has no runs - add a run to preserve paragraph formatting
                    run = first_para.add_run()
                    run.text = processed_text

                # Remove additional paragraphs to keep cell clean
                p_elements = text_frame._txBody.findall(qn('a:p'))
                for p_elem in p_elements[1:]:
                    text_frame._txBody.remove(p_elem)
            else:
                # No paragraphs at all (shouldn't happen, but fallback)
                cell.text = processed_text

            state.is_modified = True
            return f"Successfully modified cell at row {row}, column {column} in table '{table_shape.name}'"
        except Exception as e:
            return f"Error modifying table cell: {str(e)}"


def _parse_color(hex_color: str) -> RGBColor:
    """Parse a hex color string to RGBColor.

    Args:
        hex_color: Color as hex string (e.g., "#FF0000" or "FF0000")

    Returns:
        RGBColor object, or None if parsing fails
    """
    try:
        hex_color = hex_color.lstrip('#')
        if len(hex_color) == 6:
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return RGBColor(r, g, b)
        else:
            print(f"Warning: Invalid hex color '{hex_color}' - expected 6 characters")
    except ValueError as e:
        print(f"Warning: Could not parse hex color '{hex_color}': {e}")
    return None
