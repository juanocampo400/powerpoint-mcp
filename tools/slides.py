"""
Slide management tools: add, delete, duplicate, reorder slides, and get slide details.
"""
from copy import deepcopy
from pptx.util import Inches, Emu
from pptx.oxml.ns import qn


# Reverse mappings for list format detection
NUMBERED_TYPE_NAMES = {
    "arabicPeriod": "numbered (1. 2. 3.)",
    "arabicParenR": "numbered (1) 2) 3))",
    "romanLcPeriod": "roman (i. ii. iii.)",
    "romanUcPeriod": "roman (I. II. III.)",
    "alphaLcPeriod": "letter (a. b. c.)",
    "alphaUcPeriod": "letter (A. B. C.)",
}

BULLET_CHAR_NAMES = {
    "\u2022": "bullet",
    "\u2013": "dash",
    "\u2192": "arrow",
    "\u2713": "check",
    "\u25A0": "square",
    "\u25CF": "circle",
    "\u25C6": "diamond",
    "\u2605": "star",
}


def register_slide_tools(mcp, state):
    """Register slide management tools with the MCP server."""

    @mcp.tool()
    def manage_slide(
        action: str,
        slide_number: int = None,
        target_position: int = None,
        layout_index: int = 6
    ) -> str:
        """
        Manage slides: add, delete, duplicate, or move slides.

        Args:
            action: Action to perform - "add", "delete", "duplicate", or "move"
            slide_number: The slide number to operate on (1-based index). Required for delete, duplicate, move.
            target_position: For 'move' - where to move the slide. For 'add' - position to insert (optional, defaults to end).
            layout_index: For 'add' - the slide layout index to use (default 6 = blank slide).
                         Common layouts: 0=Title, 1=Title+Content, 5=Blank, 6=Blank

        Actions:
            - "add": Add a new slide (optional target_position, optional layout_index)
            - "delete": Delete a slide (requires slide_number)
            - "duplicate": Duplicate a slide (requires slide_number, optional target_position)
            - "move": Move a slide to a new position (requires slide_number and target_position)

        Returns:
            Success message with slide details, or error message
        """
        if state.presentation is None:
            return "Error: No presentation is currently open"

        prs = state.presentation
        action = action.lower().strip()
        total_slides = len(prs.slides)

        if action == "add":
            try:
                # Get layout
                if layout_index >= len(prs.slide_layouts):
                    layout_index = 6 if len(prs.slide_layouts) > 6 else len(prs.slide_layouts) - 1

                layout = prs.slide_layouts[layout_index]
                slide = prs.slides.add_slide(layout)
                state.is_modified = True

                new_position = len(prs.slides)

                # Move to target position if specified
                if target_position is not None and 1 <= target_position <= new_position:
                    # Move the newly added slide (currently at end) to target position
                    _move_slide(prs, new_position, target_position)
                    new_position = target_position

                return f"Successfully added new slide at position {new_position}\nTotal slides: {len(prs.slides)}"
            except Exception as e:
                return f"Error adding slide: {str(e)}"

        elif action == "delete":
            if slide_number is None:
                return "Error: slide_number is required for 'delete' action"
            if slide_number < 1 or slide_number > total_slides:
                return f"Error: slide_number {slide_number} is out of range (1-{total_slides})"

            try:
                # Get the slide's rId and remove it
                slide_id = prs.slides._sldIdLst[slide_number - 1].rId
                prs.part.drop_rel(slide_id)
                del prs.slides._sldIdLst[slide_number - 1]
                state.is_modified = True
                return f"Successfully deleted slide {slide_number}\nTotal slides: {len(prs.slides)}"
            except Exception as e:
                return f"Error deleting slide: {str(e)}"

        elif action == "duplicate":
            if slide_number is None:
                return "Error: slide_number is required for 'duplicate' action"
            if slide_number < 1 or slide_number > total_slides:
                return f"Error: slide_number {slide_number} is out of range (1-{total_slides})"

            try:
                # Get the source slide
                source_slide = prs.slides[slide_number - 1]

                # Add new slide with same layout
                new_slide = prs.slides.add_slide(source_slide.slide_layout)

                # Copy shapes from source to new slide
                for shape in source_slide.shapes:
                    _copy_shape(shape, new_slide)

                state.is_modified = True
                new_position = len(prs.slides)

                # Move to target position if specified
                if target_position is not None and 1 <= target_position <= new_position:
                    _move_slide(prs, new_position, target_position)
                    new_position = target_position

                return f"Successfully duplicated slide {slide_number} to position {new_position}\nTotal slides: {len(prs.slides)}"
            except Exception as e:
                return f"Error duplicating slide: {str(e)}"

        elif action == "move":
            if slide_number is None:
                return "Error: slide_number is required for 'move' action"
            if target_position is None:
                return "Error: target_position is required for 'move' action"
            if slide_number < 1 or slide_number > total_slides:
                return f"Error: slide_number {slide_number} is out of range (1-{total_slides})"
            if target_position < 1 or target_position > total_slides:
                return f"Error: target_position {target_position} is out of range (1-{total_slides})"

            if slide_number == target_position:
                return f"Slide {slide_number} is already at position {target_position}"

            try:
                _move_slide(prs, slide_number, target_position)
                state.is_modified = True
                return f"Successfully moved slide from position {slide_number} to position {target_position}"
            except Exception as e:
                return f"Error moving slide: {str(e)}"

        else:
            return f"Error: Unknown action '{action}'. Valid actions: add, delete, duplicate, move"

    @mcp.tool()
    def get_slide_snapshot(slide_number: int) -> str:
        """
        Get detailed information about a specific slide including all shapes and their properties.

        Args:
            slide_number: The slide number to inspect (1-based index)

        Returns:
            Detailed slide information including all shapes, their types, positions, and content
        """
        if state.presentation is None:
            return "Error: No presentation is currently open"

        prs = state.presentation
        total_slides = len(prs.slides)

        if slide_number < 1 or slide_number > total_slides:
            return f"Error: slide_number {slide_number} is out of range (1-{total_slides})"

        slide = prs.slides[slide_number - 1]
        info = [f"=== Slide {slide_number} of {total_slides} ==="]

        # Slide layout info
        try:
            layout_name = slide.slide_layout.name
            info.append(f"Layout: {layout_name}")
        except (AttributeError, KeyError):
            info.append("Layout: Unknown")

        info.append(f"Shape count: {len(slide.shapes)}")
        info.append("")

        # List all shapes with details
        info.append("=== Shapes ===")
        for i, shape in enumerate(slide.shapes, 1):
            shape_info = [f"\n[Shape {i}] ID: {shape.shape_id}"]
            shape_info.append(f"  Name: {shape.name}")
            shape_info.append(f"  Type: {shape.shape_type}")

            # Position and size
            shape_info.append(f"  Position: ({shape.left.inches:.2f}\", {shape.top.inches:.2f}\")")
            shape_info.append(f"  Size: {shape.width.inches:.2f}\" x {shape.height.inches:.2f}\"")

            # Text content if available
            if shape.has_text_frame:
                text = shape.text_frame.text
                if text:
                    # Truncate long text
                    display_text = text[:100] + "..." if len(text) > 100 else text
                    display_text = display_text.replace("\n", "\\n")
                    shape_info.append(f"  Text: \"{display_text}\"")

                    # Detect and display list formatting
                    list_info = _format_list_info(shape.text_frame.paragraphs)
                    if list_info:
                        shape_info.append(f"  List format: {list_info}")

            # Table info if it's a table
            if shape.has_table:
                table = shape.table
                shape_info.append(f"  Table: {len(table.rows)} rows x {len(table.columns)} columns")

            # Chart info if it's a chart
            if shape.has_chart:
                shape_info.append(f"  Chart: {shape.chart.chart_type}")

            # Icon placeholder hint
            if _is_icon_placeholder(shape):
                shape_info.append(f"  [Icon placeholder - replace with: insert_icon(slide_number={slide_number}, icon_name=\"...\", replace_shape_id={shape.shape_id})]")

            info.extend(shape_info)

        if len(slide.shapes) == 0:
            info.append("  (No shapes on this slide)")

        return "\n".join(info)


def _is_icon_placeholder(shape):
    """Check if a shape appears to be an icon placeholder.

    Detection criteria:
    - Both dimensions under 1.5 inches
    - Roughly square (within 20% tolerance)
    """
    width = shape.width.inches
    height = shape.height.inches

    # Must be small (under 1.5" on both sides)
    if width > 1.5 or height > 1.5:
        return False

    # Must be non-zero
    if max(width, height) == 0:
        return False

    # Must be roughly square (within 20% tolerance)
    ratio = min(width, height) / max(width, height)
    if ratio < 0.8:
        return False

    return True


def _detect_list_format(paragraph):
    """Detect list formatting for a single paragraph.

    Returns tuple of (format_type, format_detail, display_string) or (None, None, None).
    - format_type: 'bullet', 'numbered', or None
    - format_detail: the specific type (e.g., 'arabicPeriod', 'bullet')
    - display_string: user-friendly display (e.g., 'numbered (1. 2. 3.)')
    """
    try:
        pPr = paragraph._p.find(qn('a:pPr'))
        if pPr is None:
            return (None, None, None)

        # Check for explicitly disabled bullets
        buNone = pPr.find(qn('a:buNone'))
        if buNone is not None:
            return (None, None, None)

        # Check for auto-numbering
        buAutoNum = pPr.find(qn('a:buAutoNum'))
        if buAutoNum is not None:
            num_type = buAutoNum.get('type')
            display = NUMBERED_TYPE_NAMES.get(num_type, f"numbered ({num_type})")
            return ('numbered', num_type, display)

        # Check for character bullets
        buChar = pPr.find(qn('a:buChar'))
        if buChar is not None:
            char = buChar.get('char')
            char_name = BULLET_CHAR_NAMES.get(char, f"custom ({char})")
            return ('bullet', char, char_name)

        return (None, None, None)
    except Exception:
        return (None, None, None)


def _format_list_info(paragraphs):
    """Analyze all paragraphs and return a summary of list formatting.

    Returns a string describing the list format, or None if no list formatting.
    """
    try:
        formats = []
        for para in paragraphs:
            fmt_type, fmt_detail, fmt_display = _detect_list_format(para)
            if fmt_type is not None:
                formats.append((fmt_type, fmt_detail, fmt_display))

        if not formats:
            return None

        # Check if all formats are the same
        first_type, first_detail, first_display = formats[0]
        all_same = all(f[0] == first_type and f[1] == first_detail for f in formats)

        if all_same:
            return first_display
        else:
            # Mixed formatting
            types_present = set(f[0] for f in formats)
            if types_present == {'bullet', 'numbered'}:
                return "mixed (bullets and numbered)"
            elif 'bullet' in types_present:
                return "mixed bullets"
            elif 'numbered' in types_present:
                return "mixed numbered"
            else:
                return "mixed"
    except Exception:
        return None


def _move_slide(prs, from_pos: int, to_pos: int):
    """Move a slide from one position to another (1-based indices)."""
    slides = prs.slides._sldIdLst
    slide = slides[from_pos - 1]
    slides.remove(slide)
    slides.insert(to_pos - 1, slide)


def _copy_shape(shape, target_slide):
    """Copy a shape to a target slide.

    Handles text boxes, images, tables, and basic shapes.
    Charts and grouped shapes have limited support.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu
    import tempfile
    import os

    shape_type = shape.shape_type

    # Handle pictures/images
    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            # Extract image to temp file and re-add
            image = shape.image
            fd, temp_path = tempfile.mkstemp(suffix=f'.{image.ext}')
            os.close(fd)
            try:
                with open(temp_path, 'wb') as f:
                    f.write(image.blob)
                target_slide.shapes.add_picture(
                    temp_path,
                    shape.left, shape.top,
                    shape.width, shape.height
                )
            finally:
                os.unlink(temp_path)
        except Exception as e:
            print(f"Warning: Could not copy picture: {e}")
        return

    # Handle tables
    if shape.has_table:
        try:
            table = shape.table
            rows = len(table.rows)
            cols = len(table.columns)

            # Add new table
            new_table_shape = target_slide.shapes.add_table(
                rows, cols,
                shape.left, shape.top,
                shape.width, shape.height
            )
            new_table = new_table_shape.table

            # Copy cell content
            for r in range(rows):
                for c in range(cols):
                    src_cell = table.cell(r, c)
                    dst_cell = new_table.cell(r, c)
                    dst_cell.text = src_cell.text
        except Exception as e:
            print(f"Warning: Could not copy table: {e}")
        return

    # Handle charts (limited - just note that it exists)
    if shape.has_chart:
        print(f"Warning: Charts cannot be fully copied. Skipping chart shape.")
        return

    # Handle grouped shapes (limited support)
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        print(f"Warning: Grouped shapes cannot be fully copied. Skipping group.")
        return

    # Handle text frames (text boxes, placeholders)
    if shape.has_text_frame:
        try:
            new_shape = target_slide.shapes.add_textbox(
                shape.left, shape.top, shape.width, shape.height
            )
            # Copy text with formatting including bullets and paragraph properties
            for i, para in enumerate(shape.text_frame.paragraphs):
                if i == 0:
                    new_para = new_shape.text_frame.paragraphs[0]
                else:
                    new_para = new_shape.text_frame.add_paragraph()

                # Copy paragraph-level properties
                new_para.alignment = para.alignment
                new_para.level = para.level

                # Copy paragraph properties XML (includes bullet formatting)
                src_pPr = para._p.find(qn('a:pPr'))
                if src_pPr is not None:
                    # Remove existing pPr if present
                    existing_pPr = new_para._p.find(qn('a:pPr'))
                    if existing_pPr is not None:
                        new_para._p.remove(existing_pPr)
                    # Insert copied pPr at the beginning
                    new_para._p.insert(0, deepcopy(src_pPr))

                # Copy runs with their text and formatting
                for j, run in enumerate(para.runs):
                    if j == 0 and len(new_para.runs) > 0:
                        new_run = new_para.runs[0]
                        new_run.text = run.text
                    else:
                        new_run = new_para.add_run()
                        new_run.text = run.text
                    try:
                        if run.font.bold is not None:
                            new_run.font.bold = run.font.bold
                        if run.font.italic is not None:
                            new_run.font.italic = run.font.italic
                        if run.font.size is not None:
                            new_run.font.size = run.font.size
                        if run.font.name is not None:
                            new_run.font.name = run.font.name
                        # Copy font color if set
                        try:
                            if run.font.color.rgb is not None:
                                new_run.font.color.rgb = run.font.color.rgb
                        except AttributeError:
                            pass
                    except AttributeError:
                        pass

                # Handle case where paragraph has text but no runs (direct text)
                if not para.runs and para.text:
                    new_run = new_para.add_run()
                    new_run.text = para.text

        except Exception as e:
            print(f"Warning: Could not copy text shape: {e}")
        return

    # Handle auto shapes (rectangles, ovals, etc.)
    if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        try:
            # Get the auto shape type
            auto_shape_type = shape.auto_shape_type
            new_shape = target_slide.shapes.add_shape(
                auto_shape_type,
                shape.left, shape.top,
                shape.width, shape.height
            )
            # Copy fill if possible
            try:
                if shape.fill.type is not None:
                    new_shape.fill.solid()
                    if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color:
                        new_shape.fill.fore_color.rgb = shape.fill.fore_color.rgb
            except (AttributeError, TypeError):
                pass
            # Copy text if present
            if shape.has_text_frame and shape.text:
                new_shape.text = shape.text
        except Exception as e:
            print(f"Warning: Could not copy auto shape: {e}")
        return

    # Fallback for other shape types
    print(f"Warning: Shape type {shape_type} not fully supported for copying.")
