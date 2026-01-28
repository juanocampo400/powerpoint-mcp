"""
Escape hatch: execute arbitrary python-pptx code for advanced operations.
"""
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
import traceback


def register_evaluate_tools(mcp, state):
    """Register the evaluate_code escape hatch tool."""

    @mcp.tool()
    def evaluate_code(code: str, description: str = None) -> str:
        """
        Execute arbitrary Python code with access to python-pptx for advanced operations.

        This is an escape hatch for operations not covered by other tools.

        IMPORTANT: Before using this tool, verify standard tools cannot
        accomplish the task. Prefer find_and_replace (preserves formatting),
        modify_shape, or add_table over writing raw python-pptx code.

        FORMATTING WARNING: Direct .text assignments destroy all formatting:
            cell.text = "value"                # WRONG - destroys font, size, color
            paragraph.runs[0].text = "value"    # RIGHT - preserves run formatting

        Available in execution context:
            - prs: The current Presentation object
            - state: The PresentationState object
            - Inches, Pt, Emu: Unit helpers from pptx.util
            - MSO_SHAPE: Shape type constants
            - XL_CHART_TYPE: Chart type constants
            - PP_ALIGN, MSO_ANCHOR: Text alignment constants
            - RGBColor: Color helper
            - CategoryChartData: For creating charts

        Args:
            code: Python code to execute. Has access to the 'prs' variable (current presentation).
            description: Human-readable description of what this code does (for logging)

        Returns:
            Execution result or error message

        Example:
            evaluate_code('''
            # Access first slide
            slide = prs.slides[0]

            # Add a custom shape
            from pptx.enum.shapes import MSO_SHAPE
            shape = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                Inches(1), Inches(1),
                Inches(2), Inches(1)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0, 112, 192)
            result = f"Added chevron shape: {shape.shape_id}"
            ''', description="Add a blue chevron shape")
        """
        if state.presentation is None:
            return "Error: No presentation is currently open"

        # Build execution context
        exec_globals = {
            # The presentation
            'prs': state.presentation,
            'state': state,

            # Common imports pre-loaded
            'Inches': Inches,
            'Pt': Pt,
            'Emu': Emu,
            'MSO_SHAPE': MSO_SHAPE,
            'XL_CHART_TYPE': XL_CHART_TYPE,
            'PP_ALIGN': PP_ALIGN,
            'MSO_ANCHOR': MSO_ANCHOR,
            'RGBColor': RGBColor,
            'CategoryChartData': CategoryChartData,

            # Result variable
            'result': None
        }

        try:
            # Execute the code
            exec(code, exec_globals)

            # Mark as modified since we don't know what the code did
            state.is_modified = True

            # Return result if set, otherwise generic success
            result = exec_globals.get('result')
            if result is not None:
                return f"Code executed successfully.\nResult: {result}"
            else:
                desc_msg = f" ({description})" if description else ""
                return f"Code executed successfully{desc_msg}."

        except Exception as e:
            error_trace = traceback.format_exc()
            return f"Error executing code:\n{error_trace}"
