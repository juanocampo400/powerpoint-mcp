"""
Icon tools: insert Phosphor SVG icons into presentations with recolorable support.
"""
import logging
import os
import tempfile
from pathlib import Path

from pptx.util import Inches

from .svg_embed import SVGEmbedder, make_svg_recolorable, generate_png_fallback
from .shape_utils import get_shape_and_geometry, delete_shape

logger = logging.getLogger(__name__)

# Get the icons directory relative to this file
ICONS_DIR = Path(__file__).parent.parent / "icons" / "phosphor"

# Common Phosphor icons organized by category (fill variant, but user specifies without -fill)
AVAILABLE_ICONS = {
    "arrows": [
        "arrow-up", "arrow-down", "arrow-left", "arrow-right",
        "arrow-circle-up", "arrow-circle-down", "arrow-circle-left", "arrow-circle-right",
        "caret-up", "caret-down", "caret-left", "caret-right",
        "caret-circle-up", "caret-circle-down", "caret-circle-left", "caret-circle-right",
        "arrows-clockwise", "arrows-counter-clockwise", "arrow-clockwise", "arrow-counter-clockwise",
        "arrows-horizontal", "arrows-vertical", "arrows-in", "arrows-out"
    ],
    "communication": [
        "envelope", "envelope-simple", "envelope-open", "phone", "phone-call",
        "chat", "chat-circle", "chat-dots", "chats", "chat-text",
        "bell", "bell-ringing", "bell-simple", "megaphone", "broadcast"
    ],
    "data": [
        "chart-bar", "chart-bar-horizontal", "chart-line", "chart-line-up", "chart-line-down",
        "chart-pie", "chart-pie-slice", "chart-donut", "chart-polar", "chart-scatter",
        "graph", "pulse", "heartbeat", "trend-up", "trend-down"
    ],
    "files": [
        "file", "file-text", "file-pdf", "file-doc", "file-xls", "file-ppt",
        "file-code", "file-image", "file-video", "file-audio", "file-zip",
        "folder", "folder-open", "folder-plus", "folders",
        "clipboard", "clipboard-text", "copy", "download", "upload"
    ],
    "general": [
        "house", "house-simple", "gear", "gear-six", "sliders", "sliders-horizontal",
        "magnifying-glass", "funnel", "funnel-simple",
        "list", "list-bullets", "list-numbers", "grid-four", "squares-four",
        "dots-three", "dots-three-vertical", "dots-nine"
    ],
    "media": [
        "image", "images", "camera", "video", "video-camera", "film-strip",
        "microphone", "microphone-slash", "speaker-high", "speaker-low", "speaker-none",
        "play", "pause", "stop", "skip-forward", "skip-back", "rewind", "fast-forward"
    ],
    "people": [
        "user", "user-circle", "user-plus", "user-minus", "user-check", "user-gear",
        "users", "users-three", "users-four", "person", "person-simple"
    ],
    "status": [
        "check", "check-circle", "check-square", "check-fat",
        "x", "x-circle", "x-square",
        "warning", "warning-circle", "warning-diamond", "warning-octagon",
        "info", "question", "prohibit", "seal-check", "seal-warning"
    ],
    "time": [
        "clock", "clock-countdown", "clock-clockwise", "clock-counter-clockwise",
        "calendar", "calendar-blank", "calendar-check", "calendar-plus", "calendar-x",
        "timer", "hourglass", "hourglass-simple", "alarm"
    ],
    "weather": [
        "sun", "sun-dim", "moon", "moon-stars", "cloud", "cloud-sun", "cloud-moon",
        "cloud-rain", "cloud-snow", "cloud-lightning", "snowflake", "thermometer",
        "wind", "rainbow", "umbrella", "drop"
    ],
    "business": [
        "briefcase", "building", "building-office", "buildings", "bank",
        "currency-dollar", "currency-eur", "money", "wallet", "credit-card",
        "shopping-cart", "shopping-bag", "storefront", "package", "truck"
    ],
    "misc": [
        "star", "heart", "bookmark", "flag", "tag", "link", "link-break",
        "lock", "lock-open", "key", "shield", "shield-check",
        "eye", "eye-slash", "pencil", "trash", "plus", "minus",
        "lightning", "fire", "globe", "map-pin", "target", "trophy", "medal", "gift",
        "lightbulb", "rocket", "puzzle-piece", "magic-wand"
    ]
}

# Pre-compute flattened list for suggestions
ALL_CURATED_ICONS = [icon for icons in AVAILABLE_ICONS.values() for icon in icons]

# Check cairosvg availability at module load
try:
    import cairosvg
    CAIROSVG_AVAILABLE = True
except (ImportError, OSError):
    CAIROSVG_AVAILABLE = False
    logger.warning("cairosvg not available - icon insertion will not work. Install cairo: brew install cairo (macOS) or apt install libcairo2-dev (Linux)")


def _get_icon_svg_path(icon_name: str) -> Path:
    """Get the path to an icon's SVG file.

    Phosphor icons are stored with -fill suffix internally,
    but users specify the icon name without the suffix.

    Args:
        icon_name: Icon name without -fill suffix (e.g., "check-circle")

    Returns:
        Path to the SVG file
    """
    return ICONS_DIR / f"{icon_name}-fill.svg"


def _load_icon_svg(icon_name: str) -> str:
    """Load SVG content for an icon.

    Args:
        icon_name: Icon name without -fill suffix

    Returns:
        SVG content as string

    Raises:
        FileNotFoundError: If icon doesn't exist
    """
    svg_path = _get_icon_svg_path(icon_name)
    if not svg_path.exists():
        raise FileNotFoundError(f"Icon '{icon_name}' not found at {svg_path}")

    with open(svg_path, 'r', encoding='utf-8') as f:
        return f.read()


def register_icon_tools(mcp, state):
    """Register icon tools with the MCP server."""

    @mcp.tool()
    def list_icons() -> str:
        """
        List all available Phosphor icons organized by category.

        Returns:
            Formatted list of available icons grouped by category
        """
        lines = ["=== Phosphor Icons (Fill Variant) ==="]
        lines.append("Over 1,500 icons available. Common icons listed below.")
        lines.append("Browse all: https://phosphoricons.com/\n")

        for category, icons in AVAILABLE_ICONS.items():
            lines.append(f"\n{category.upper()}:")
            lines.append(f"  {', '.join(icons)}")

        lines.append("\n\nUsage: insert_icon(slide_number=1, icon_name='check-circle')")
        lines.append("The 'color' parameter sets the initial icon color (default #333333).")
        lines.append("Users can change colors in PowerPoint via Graphics Format > Graphics Fill.")

        return "\n".join(lines)

    @mcp.tool()
    def insert_icon(
        slide_number: int,
        icon_name: str,
        left: float = None,
        top: float = None,
        size: float = None,
        color: str = "#333333",
        replace_shape_id: int = None,
        replace_shape_name: str = None
    ) -> str:
        """
        Insert a Phosphor icon into a slide.

        Args:
            slide_number: Target slide number (1-based)
            icon_name: Name of the Phosphor icon (e.g., "check-circle", "user", "star")
                      Use list_icons() to see all available icons
            left: Left position in inches (default 1.0, or inherited from replaced shape)
            top: Top position in inches (default 1.0, or inherited from replaced shape)
            size: Icon size in inches (default 1.0, or inherited from replaced shape) - icons are square
            color: Icon color as hex code (default "#333333"). The icon displays with
                   this color and remains recolorable in PowerPoint via Graphics Fill.
            replace_shape_id: ID of shape to replace (icon inherits its position/size)
            replace_shape_name: Name of shape to replace (alternative to replace_shape_id)

        Returns:
            Success message with shape ID, or error message

        Note:
            When using replace_shape_id or replace_shape_name, the icon will inherit
            the position and size of the replaced shape (using min of width/height since
            icons are square). You can still override individual properties by providing
            explicit values for left, top, or size.
        """
        if not CAIROSVG_AVAILABLE:
            return "Error: cairosvg is required for icon insertion. Run: pip install cairosvg"

        if state.presentation is None:
            return "Error: No presentation is currently open"

        prs = state.presentation
        if slide_number < 1 or slide_number > len(prs.slides):
            return f"Error: slide_number {slide_number} is out of range (1-{len(prs.slides)})"

        slide = prs.slides[slide_number - 1]

        # Validate icon exists BEFORE modifying any shapes
        svg_path = _get_icon_svg_path(icon_name)
        if not svg_path.exists():
            similar = [i for i in ALL_CURATED_ICONS if icon_name.lower() in i.lower()]
            suggestion = f" Similar icons: {', '.join(similar[:5])}" if similar else ""
            return f"Error: Icon '{icon_name}' not found.{suggestion}\nBrowse all icons at: https://phosphoricons.com/"

        # Handle replace mode (now safe to delete since icon exists)
        if replace_shape_id is not None or replace_shape_name is not None:
            shape_result, geometry_or_error = get_shape_and_geometry(
                slide, replace_shape_id, replace_shape_name
            )
            if shape_result is None:
                return f"Error: {geometry_or_error}"

            # Use placeholder geometry as defaults
            if left is None:
                left = geometry_or_error['left']
            if top is None:
                top = geometry_or_error['top']
            if size is None:
                # Icons are square, use smaller dimension
                size = min(geometry_or_error['width'], geometry_or_error['height'])

            # Delete the placeholder (safe: icon existence already validated)
            delete_shape(shape_result)

        # Apply defaults for non-replace mode
        if left is None:
            left = 1.0
        if top is None:
            top = 1.0
        if size is None:
            size = 1.0

        try:
            # Load and process SVG
            svg_content = _load_icon_svg(icon_name)

            # Make SVG recolorable (strip color attributes, apply fill color)
            recolorable_svg = make_svg_recolorable(svg_content, fill_color=color)

            # Generate PNG fallback with specified color
            size_px = max(96, int(size * 96))  # At least 96px, or scale to size
            png_bytes = generate_png_fallback(svg_content, color, size_px)

            # Embed with dual SVG+PNG format
            embedder = SVGEmbedder()
            shape_id = embedder.embed_recolorable_icon(
                slide=slide,
                svg_content=recolorable_svg,
                png_bytes=png_bytes,
                left_inches=left,
                top_inches=top,
                size_inches=size,
                icon_name=icon_name
            )

            state.is_modified = True
            color_msg = f" with color {color}" if color != "#333333" else ""
            return (
                f"Successfully added '{icon_name}' icon on slide {slide_number}{color_msg}\n"
                f"Shape ID: {shape_id}\n"
                f"Size: {size}\" x {size}\"\n"
                f"Recolorable: Yes (use Graphics Format > Graphics Fill in PowerPoint)"
            )

        except Exception as e:
            logger.exception(f"Error inserting icon '{icon_name}'")
            return f"Error inserting icon: {str(e)}"
