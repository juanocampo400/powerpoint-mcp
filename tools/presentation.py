"""
Presentation management tools: open, create, save, close presentations.
"""
from pptx import Presentation
from pathlib import Path
import os

from .svg_embed import ensure_svg_content_type


def register_presentation_tools(mcp, state):
    """Register presentation management tools with the MCP server."""

    @mcp.tool()
    def manage_presentation(
        action: str,
        file_path: str = None,
        save_path: str = None
    ) -> str:
        """
        Manage PowerPoint presentations: open, create, save, save_as, or close.

        Args:
            action: Action to perform - "open", "create", "save", "save_as", or "close"
            file_path: Path for open/create operations (required for open)
            save_path: New path for save_as operation (required for save_as)

        Actions:
            - "open": Opens an existing presentation (requires file_path)
            - "create": Creates a new blank presentation (optional file_path to set default save location)
            - "save": Saves the current presentation to its original location
            - "save_as": Saves the current presentation to a new location (requires save_path)
            - "close": Closes the current presentation without saving

        Returns:
            Success message with operation details, or error message
        """
        action = action.lower().strip()

        if action == "open":
            if not file_path:
                return "Error: file_path is required for 'open' action"

            # Normalize path
            file_path = os.path.normpath(os.path.expanduser(file_path))

            if not os.path.exists(file_path):
                return f"Error: File not found: {file_path}"

            try:
                state.presentation = Presentation(file_path)
                state.file_path = file_path
                state.is_modified = False
                slide_count = len(state.presentation.slides)
                return f"Successfully opened presentation: {file_path}\nSlide count: {slide_count}"
            except Exception as e:
                return f"Error opening presentation: {str(e)}"

        elif action == "create":
            try:
                state.presentation = Presentation()
                state.is_modified = True
                if file_path:
                    state.file_path = os.path.normpath(os.path.expanduser(file_path))
                else:
                    state.file_path = None
                return f"Successfully created new blank presentation" + (f"\nDefault save path: {state.file_path}" if state.file_path else "\nNo save path set - use save_as to save")
            except Exception as e:
                return f"Error creating presentation: {str(e)}"

        elif action == "save":
            if state.presentation is None:
                return "Error: No presentation is currently open"
            if not state.file_path:
                return "Error: No file path set. Use 'save_as' with a save_path instead"

            try:
                state.presentation.save(state.file_path)
                # Ensure SVG content type is registered for recolorable icons
                ensure_svg_content_type(state.file_path)
                state.is_modified = False
                return f"Successfully saved presentation to: {state.file_path}"
            except Exception as e:
                return f"Error saving presentation: {str(e)}"

        elif action == "save_as":
            if state.presentation is None:
                return "Error: No presentation is currently open"
            if not save_path:
                return "Error: save_path is required for 'save_as' action"

            save_path = os.path.normpath(os.path.expanduser(save_path))

            # Ensure directory exists
            save_dir = os.path.dirname(save_path)
            if save_dir and not os.path.exists(save_dir):
                os.makedirs(save_dir, exist_ok=True)

            try:
                state.presentation.save(save_path)
                # Ensure SVG content type is registered for recolorable icons
                ensure_svg_content_type(save_path)
                state.file_path = save_path
                state.is_modified = False
                return f"Successfully saved presentation to: {save_path}"
            except Exception as e:
                return f"Error saving presentation: {str(e)}"

        elif action == "close":
            if state.presentation is None:
                return "No presentation is currently open"

            file_info = state.file_path or "unsaved presentation"
            was_modified = state.is_modified
            state.reset()

            msg = f"Closed presentation: {file_info}"
            if was_modified:
                msg += "\nWarning: Unsaved changes were discarded"
            return msg

        else:
            return f"Error: Unknown action '{action}'. Valid actions: open, create, save, save_as, close"

    @mcp.tool()
    def get_presentation_info() -> str:
        """
        Get information about the currently open presentation.

        Returns:
            Presentation details including file path, slide count, dimensions, and modification status.
        """
        if state.presentation is None:
            return "No presentation is currently open. Use manage_presentation(action='open') or manage_presentation(action='create') first."

        prs = state.presentation
        slide_count = len(prs.slides)

        # Get slide dimensions
        width_inches = prs.slide_width.inches
        height_inches = prs.slide_height.inches

        info = [
            "=== Presentation Info ===",
            f"File path: {state.file_path or 'Not saved yet'}",
            f"Slide count: {slide_count}",
            f"Slide dimensions: {width_inches:.2f}\" x {height_inches:.2f}\"",
            f"Modified: {'Yes' if state.is_modified else 'No'}",
        ]

        if slide_count > 0:
            info.append("\n=== Slides Overview ===")
            for i, slide in enumerate(prs.slides, 1):
                shape_count = len(slide.shapes)
                # Try to get title
                title = "No title"
                if slide.shapes.title:
                    title = slide.shapes.title.text[:50] or "Empty title"
                info.append(f"  Slide {i}: {shape_count} shapes - \"{title}\"")

        return "\n".join(info)
