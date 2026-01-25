"""
Content tools: add textboxes, images, shapes, tables, and charts.
"""
from lxml import etree
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.shapes.placeholder import PicturePlaceholder
import json
import os

from .shape_utils import get_shape_and_geometry, delete_shape


# Valid fit modes for add_image
VALID_FIT_MODES = ["stretch", "fill", "fit"]

# Bullet type mappings (shared with modify.py)
BULLET_CHARS = {
    "bullet": "\u2022",      # •
    "dash": "\u2013",        # –
    "arrow": "\u2192",       # →
    "check": "\u2713",       # ✓
    "square": "\u25A0",      # ■
    "circle": "\u25CF",      # ●
    "diamond": "\u25C6",     # ◆
    "star": "\u2605",        # ★
}

NUMBERED_TYPES = {
    "number": "arabicPeriod",       # 1. 2. 3.
    "number_paren": "arabicParenR", # 1) 2) 3)
    "roman": "romanLcPeriod",       # i. ii. iii.
    "roman_upper": "romanUcPeriod", # I. II. III.
    "letter": "alphaLcPeriod",      # a. b. c.
    "letter_upper": "alphaUcPeriod", # A. B. C.
}

ALL_BULLET_TYPES = list(BULLET_CHARS.keys()) + list(NUMBERED_TYPES.keys()) + ["none"]


def _process_text_escapes(text: str) -> str:
    """Convert escape sequences like \\n to actual newlines.

    Args:
        text: Input text that may contain escape sequences

    Returns:
        Text with escape sequences converted to actual characters
    """
    if text is None:
        return None
    return text.replace("\\n", "\n").replace("\\t", "\t")


def _apply_fit_mode(picture, target_width, target_height, fit_mode, image_path):
    """Apply fit mode (fill/fit) to an image using crop values.

    Args:
        picture: The Picture shape object (already added to slide)
        target_width: Target width in Emu
        target_height: Target height in Emu
        fit_mode: "fill" (crop to fill) or "fit" (fit within bounds)
        image_path: Path to the image file (for getting dimensions)

    Returns:
        The modified picture object
    """
    from PIL import Image

    # Get natural image dimensions
    with Image.open(image_path) as img:
        img_width, img_height = img.size

    img_aspect = img_width / img_height
    target_aspect = target_width / target_height

    if fit_mode == "fill":
        # Scale to fill (cover) - image fills entire target area, overflow is cropped
        if img_aspect > target_aspect:
            # Image is wider than target - crop left/right
            # Scale by height to fill vertically
            scale = target_height / img_height
            scaled_width = img_width * scale
            # Calculate crop as proportion of scaled image
            excess = (scaled_width - target_width) / scaled_width
            crop_x = excess / 2
            crop_y = 0
        else:
            # Image is taller than target - crop top/bottom
            # Scale by width to fill horizontally
            scale = target_width / img_width
            scaled_height = img_height * scale
            # Calculate crop as proportion of scaled image
            excess = (scaled_height - target_height) / scaled_height
            crop_x = 0
            crop_y = excess / 2

        # Apply crops and set final dimensions
        picture.crop_left = crop_x
        picture.crop_right = crop_x
        picture.crop_top = crop_y
        picture.crop_bottom = crop_y
        picture.width = int(target_width)
        picture.height = int(target_height)

    elif fit_mode == "fit":
        # Scale to fit (contain) - entire image visible within target area
        if img_aspect > target_aspect:
            # Image is wider - fit to width, center vertically
            new_width = target_width
            new_height = int(target_width / img_aspect)
        else:
            # Image is taller - fit to height, center horizontally
            new_height = target_height
            new_width = int(target_height * img_aspect)

        # Center within target area
        offset_x = (target_width - new_width) / 2
        offset_y = (target_height - new_height) / 2

        picture.width = int(new_width)
        picture.height = int(new_height)
        picture.left = picture.left + int(offset_x)
        picture.top = picture.top + int(offset_y)

    return picture


def _apply_bullet_to_paragraph(paragraph, bullet_type):
    """Apply a specific bullet style to a paragraph.

    Args:
        paragraph: A python-pptx Paragraph object
        bullet_type: One of BULLET_CHARS keys, NUMBERED_TYPES keys, or "none"
    """
    p_elem = paragraph._p

    # Get or create pPr element
    pPr = p_elem.find(qn('a:pPr'))
    if pPr is None:
        pPr = etree.Element(qn('a:pPr'))
        p_elem.insert(0, pPr)

    # Remove existing bullet elements
    for elem_name in ['a:buNone', 'a:buChar', 'a:buAutoNum']:
        existing = pPr.find(qn(elem_name))
        if existing is not None:
            pPr.remove(existing)

    if bullet_type == "none":
        # Explicitly no bullets
        buNone = etree.SubElement(pPr, qn('a:buNone'))
    elif bullet_type in BULLET_CHARS:
        # Character bullet
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', BULLET_CHARS[bullet_type])
    elif bullet_type in NUMBERED_TYPES:
        # Numbered list
        buAutoNum = etree.SubElement(pPr, qn('a:buAutoNum'))
        buAutoNum.set('type', NUMBERED_TYPES[bullet_type])


def register_content_tools(mcp, state):
    """Register content creation tools with the MCP server."""

    @mcp.tool()
    def add_textbox(
        slide_number: int,
        text: str,
        left: float = 1.0,
        top: float = 1.0,
        width: float = 8.0,
        height: float = 1.0,
        font_name: str = None,
        font_size: int = None,
        font_bold: bool = False,
        font_italic: bool = False,
        font_color: str = None,
        alignment: str = "left",
        bullets: str = None
    ) -> str:
        """
        Add a textbox to a slide.

        Args:
            slide_number: Target slide number (1-based)
            text: The text content to add. Use \\n to create multiple lines/bullet points.
            left: Left position in inches (default 1.0)
            top: Top position in inches (default 1.0)
            width: Width in inches (default 8.0)
            height: Height in inches (default 1.0)
            font_name: Font family name (e.g., "Arial", "Calibri")
            font_size: Font size in points (e.g., 18, 24, 36)
            font_bold: Make text bold (default False)
            font_italic: Make text italic (default False)
            font_color: Hex color code (e.g., "#FF0000" for red)
            alignment: Text alignment - "left", "center", "right" (default "left")
            bullets: Bullet/list type. Options:
                - Character bullets: "bullet", "dash", "arrow", "check", "square", "circle", "diamond", "star"
                - Numbered lists: "number" (1. 2. 3.), "number_paren" (1) 2) 3)),
                  "roman" (i. ii. iii.), "roman_upper" (I. II. III.),
                  "letter" (a. b. c.), "letter_upper" (A. B. C.)
                - "none" to explicitly disable bullets
                - None (default) for no bullets

        Returns:
            Success message with shape ID, or error message
        """
        if state.presentation is None:
            return "Error: No presentation is currently open"

        prs = state.presentation
        if slide_number < 1 or slide_number > len(prs.slides):
            return f"Error: slide_number {slide_number} is out of range (1-{len(prs.slides)})"

        slide = prs.slides[slide_number - 1]

        # Validate bullet type if specified
        if bullets is not None and bullets not in ALL_BULLET_TYPES:
            return f"Error: Invalid bullet type '{bullets}'. Valid types: {', '.join(ALL_BULLET_TYPES)}"

        try:
            # Create textbox
            shape = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            tf = shape.text_frame
            tf.word_wrap = True
            processed_text = _process_text_escapes(text)

            # Alignment mapping
            align_map = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT
            }
            para_alignment = align_map.get(alignment.lower(), PP_ALIGN.LEFT)

            # Parse color once if specified
            parsed_color = None
            if font_color:
                parsed_color = _parse_color(font_color)

            # Split text into lines for multi-paragraph support
            lines = processed_text.split('\n') if processed_text else ['']

            for i, line in enumerate(lines):
                if i == 0:
                    # Use the first (existing) paragraph
                    p = tf.paragraphs[0]
                else:
                    # Add new paragraph for subsequent lines
                    p = tf.add_paragraph()

                # Add text via run
                run = p.add_run()
                run.text = line

                # Apply font formatting
                if font_name:
                    run.font.name = font_name
                if font_size:
                    run.font.size = Pt(font_size)
                run.font.bold = font_bold
                run.font.italic = font_italic

                if parsed_color:
                    run.font.color.rgb = parsed_color

                # Apply alignment
                p.alignment = para_alignment

                # Apply bullet formatting if specified
                if bullets is not None:
                    _apply_bullet_to_paragraph(p, bullets)

            state.is_modified = True
            return f"Successfully added textbox on slide {slide_number}\nShape ID: {shape.shape_id}\nName: {shape.name}"
        except Exception as e:
            return f"Error adding textbox: {str(e)}"

    @mcp.tool()
    def add_image(
        slide_number: int,
        image_path: str,
        left: float = None,
        top: float = None,
        width: float = None,
        height: float = None,
        fit_mode: str = None,
        replace_shape_id: int = None,
        replace_shape_name: str = None
    ) -> str:
        """
        Add an image to a slide.

        Args:
            slide_number: Target slide number (1-based)
            image_path: Path to the image file (supports PNG, JPG, GIF, BMP, etc.)
            left: Left position in inches (default 1.0, or inherited from replaced shape)
            top: Top position in inches (default 1.0, or inherited from replaced shape)
            width: Width in inches (optional - maintains aspect ratio if only one dimension specified)
            height: Height in inches (optional - maintains aspect ratio if only one dimension specified)
            fit_mode: How to fit the image when both width and height are specified:
                - "fill": Scale to fill the area completely, cropping overflow (preserves aspect ratio)
                - "fit": Scale to fit within the area, may have empty space (preserves aspect ratio)
                - "stretch": Stretch to exact dimensions (distorts if aspect ratios differ)
                - None (default): "fill" when replacing a picture placeholder, "stretch" otherwise
            replace_shape_id: ID of shape to replace (image inherits its position/size)
            replace_shape_name: Name of shape to replace (alternative to replace_shape_id)

        Returns:
            Success message with shape ID, or error message

        Note:
            When replacing an empty picture placeholder (PicturePlaceholder), the image is
            automatically inserted with "fill" behavior using PowerPoint's native cropping.
            For other shapes or pre-filled placeholders, manual crop calculations are applied.
        """
        if state.presentation is None:
            return "Error: No presentation is currently open"

        prs = state.presentation
        if slide_number < 1 or slide_number > len(prs.slides):
            return f"Error: slide_number {slide_number} is out of range (1-{len(prs.slides)})"

        # Validate fit_mode
        if fit_mode is not None and fit_mode not in VALID_FIT_MODES:
            return f"Error: Invalid fit_mode '{fit_mode}'. Valid options: {', '.join(VALID_FIT_MODES)}"

        image_path = os.path.normpath(os.path.expanduser(image_path))
        if not os.path.exists(image_path):
            return f"Error: Image file not found: {image_path}"

        slide = prs.slides[slide_number - 1]

        # Track if user explicitly provided dimensions (for deciding insert_picture behavior)
        user_provided_width = width is not None
        user_provided_height = height is not None

        # Handle replace mode
        if replace_shape_id is not None or replace_shape_name is not None:
            shape_result, geometry_or_error = get_shape_and_geometry(
                slide, replace_shape_id, replace_shape_name
            )
            if shape_result is None:
                return f"Error: {geometry_or_error}"

            # Check if this is an empty PicturePlaceholder - can use native insert_picture
            is_empty_picture_placeholder = isinstance(shape_result, PicturePlaceholder)

            # Use native insert_picture if:
            # 1. It's an empty PicturePlaceholder
            # 2. User didn't override dimensions (wants placeholder's size)
            # 3. fit_mode is not "stretch" (user doesn't want distortion)
            if (is_empty_picture_placeholder and
                not user_provided_width and
                not user_provided_height and
                fit_mode != "stretch"):

                try:
                    # Use PowerPoint's native fill+crop behavior
                    picture = shape_result.insert_picture(image_path)
                    state.is_modified = True

                    # Report crop values if any
                    crop_info = ""
                    if picture.crop_top > 0 or picture.crop_bottom > 0:
                        crop_info = f"\nAuto-cropped: {picture.crop_top*100:.1f}% top, {picture.crop_bottom*100:.1f}% bottom"
                    elif picture.crop_left > 0 or picture.crop_right > 0:
                        crop_info = f"\nAuto-cropped: {picture.crop_left*100:.1f}% left, {picture.crop_right*100:.1f}% right"

                    return (f"Successfully added image on slide {slide_number}\n"
                            f"Shape ID: {picture.shape_id}\n"
                            f"Name: {picture.name}\n"
                            f"Size: {picture.width.inches:.2f}\" x {picture.height.inches:.2f}\""
                            f"{crop_info}")
                except Exception as e:
                    return f"Error inserting picture into placeholder: {str(e)}"

            # For non-PicturePlaceholder shapes or when user overrides dimensions:
            # Use placeholder geometry as defaults
            if left is None:
                left = geometry_or_error['left']
            if top is None:
                top = geometry_or_error['top']
            if width is None:
                width = geometry_or_error['width']
            if height is None:
                height = geometry_or_error['height']

            # Default to "fill" when replacing shapes (unless user specified otherwise)
            if fit_mode is None:
                fit_mode = "fill"

            # Delete the shape (will be replaced by new image)
            delete_shape(shape_result)

        # Apply defaults for non-replace mode
        if left is None:
            left = 1.0
        if top is None:
            top = 1.0

        # Default fit_mode for non-replace: stretch (preserve original behavior)
        if fit_mode is None:
            fit_mode = "stretch"

        try:
            # Determine how to add the image based on fit_mode
            if fit_mode == "stretch" or width is None or height is None:
                # Original behavior: stretch to exact dimensions or preserve aspect ratio
                width_val = Inches(width) if width else None
                height_val = Inches(height) if height else None

                shape = slide.shapes.add_picture(
                    image_path,
                    Inches(left), Inches(top),
                    width=width_val, height=height_val
                )
            else:
                # fill or fit mode: add at natural size first, then apply fit mode
                shape = slide.shapes.add_picture(
                    image_path,
                    Inches(left), Inches(top)
                )

                # Apply fit mode with crop/resize
                target_width = Inches(width)
                target_height = Inches(height)
                _apply_fit_mode(shape, target_width, target_height, fit_mode, image_path)

            state.is_modified = True

            # Build response message
            fit_info = ""
            if fit_mode == "fill" and (shape.crop_top > 0 or shape.crop_bottom > 0 or
                                        shape.crop_left > 0 or shape.crop_right > 0):
                if shape.crop_top > 0 or shape.crop_bottom > 0:
                    fit_info = f"\nCropped: {shape.crop_top*100:.1f}% top, {shape.crop_bottom*100:.1f}% bottom"
                else:
                    fit_info = f"\nCropped: {shape.crop_left*100:.1f}% left, {shape.crop_right*100:.1f}% right"
            elif fit_mode == "fit":
                fit_info = "\nFit within bounds (aspect ratio preserved)"

            return (f"Successfully added image on slide {slide_number}\n"
                    f"Shape ID: {shape.shape_id}\n"
                    f"Name: {shape.name}\n"
                    f"Size: {shape.width.inches:.2f}\" x {shape.height.inches:.2f}\""
                    f"{fit_info}")
        except Exception as e:
            return f"Error adding image: {str(e)}"

    @mcp.tool()
    def add_shape(
        slide_number: int,
        shape_type: str,
        left: float = 1.0,
        top: float = 1.0,
        width: float = 2.0,
        height: float = 2.0,
        fill_color: str = None,
        line_color: str = None,
        line_width: float = None,
        text: str = None
    ) -> str:
        """
        Add a shape to a slide.

        Args:
            slide_number: Target slide number (1-based)
            shape_type: Shape type - "rectangle", "oval", "rounded_rectangle", "triangle",
                       "right_arrow", "left_arrow", "up_arrow", "down_arrow",
                       "star", "pentagon", "hexagon", "diamond", "line"
            left: Left position in inches (default 1.0)
            top: Top position in inches (default 1.0)
            width: Width in inches (default 2.0)
            height: Height in inches (default 2.0)
            fill_color: Fill color as hex code (e.g., "#0066CC")
            line_color: Line/border color as hex code
            line_width: Line width in points
            text: Optional text to add inside the shape

        Returns:
            Success message with shape ID, or error message
        """
        if state.presentation is None:
            return "Error: No presentation is currently open"

        prs = state.presentation
        if slide_number < 1 or slide_number > len(prs.slides):
            return f"Error: slide_number {slide_number} is out of range (1-{len(prs.slides)})"

        slide = prs.slides[slide_number - 1]

        # Map shape type names to MSO_SHAPE constants
        shape_map = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "oval": MSO_SHAPE.OVAL,
            "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
            "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
            "right_arrow": MSO_SHAPE.RIGHT_ARROW,
            "left_arrow": MSO_SHAPE.LEFT_ARROW,
            "up_arrow": MSO_SHAPE.UP_ARROW,
            "down_arrow": MSO_SHAPE.DOWN_ARROW,
            "star": MSO_SHAPE.STAR_5_POINT,
            "pentagon": MSO_SHAPE.PENTAGON,
            "hexagon": MSO_SHAPE.HEXAGON,
            "diamond": MSO_SHAPE.DIAMOND,
            "line": MSO_SHAPE.LINE_INVERSE,
        }

        shape_type_lower = shape_type.lower()
        if shape_type_lower not in shape_map:
            return f"Error: Unknown shape type '{shape_type}'. Valid types: {', '.join(shape_map.keys())}"

        try:
            shape = slide.shapes.add_shape(
                shape_map[shape_type_lower],
                Inches(left), Inches(top),
                Inches(width), Inches(height)
            )

            # Apply fill color
            if fill_color:
                color = _parse_color(fill_color)
                if color:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = color

            # Apply line color
            if line_color:
                color = _parse_color(line_color)
                if color:
                    shape.line.color.rgb = color

            # Apply line width
            if line_width:
                shape.line.width = Pt(line_width)

            # Add text if specified
            if text:
                shape.text = _process_text_escapes(text)

            state.is_modified = True
            return f"Successfully added {shape_type} shape on slide {slide_number}\nShape ID: {shape.shape_id}\nName: {shape.name}"
        except Exception as e:
            return f"Error adding shape: {str(e)}"

    @mcp.tool()
    def add_table(
        slide_number: int,
        rows: int,
        cols: int,
        data: str = None,
        left: float = None,
        top: float = None,
        width: float = None,
        height: float = None,
        replace_shape_id: int = None,
        replace_shape_name: str = None
    ) -> str:
        """
        Add a table to a slide.

        Args:
            slide_number: Target slide number (1-based)
            rows: Number of rows
            cols: Number of columns
            data: JSON string of 2D array for cell data, e.g., '[["A","B"],["1","2"]]'
                  First row is typically used as header.
            left: Left position in inches (default 1.0, or inherited from replaced shape)
            top: Top position in inches (default 2.0, or inherited from replaced shape)
            width: Width in inches (default 8.0, or inherited from replaced shape)
            height: Height in inches (default 3.0, or inherited from replaced shape)
            replace_shape_id: ID of shape to replace (table inherits its position/size)
            replace_shape_name: Name of shape to replace (alternative to replace_shape_id)

        Returns:
            Success message with shape ID, or error message

        Note:
            When using replace_shape_id or replace_shape_name, the table will inherit
            the position and size of the replaced shape. You can still override individual
            properties by providing explicit values for left, top, width, or height.
        """
        if state.presentation is None:
            return "Error: No presentation is currently open"

        prs = state.presentation
        if slide_number < 1 or slide_number > len(prs.slides):
            return f"Error: slide_number {slide_number} is out of range (1-{len(prs.slides)})"

        if rows < 1 or cols < 1:
            return "Error: rows and cols must be at least 1"

        slide = prs.slides[slide_number - 1]

        # Handle replace mode
        if replace_shape_id is not None or replace_shape_name is not None:
            shape_result, geometry_or_error = get_shape_and_geometry(
                slide, replace_shape_id, replace_shape_name
            )
            if shape_result is None:
                return f"Error: {geometry_or_error}"

            # Use placeholder geometry as defaults
            if left is None:
                left = geometry_or_error['left']
            if top is None:
                top = geometry_or_error['top']
            if width is None:
                width = geometry_or_error['width']
            if height is None:
                height = geometry_or_error['height']

            # Delete the placeholder
            delete_shape(shape_result)

        # Apply defaults for non-replace mode
        if left is None:
            left = 1.0
        if top is None:
            top = 2.0
        if width is None:
            width = 8.0
        if height is None:
            height = 3.0

        try:
            # Add table
            shape = slide.shapes.add_table(
                rows, cols,
                Inches(left), Inches(top),
                Inches(width), Inches(height)
            )
            table = shape.table

            # Populate with data if provided
            if data:
                try:
                    data_array = json.loads(data)
                    for row_idx, row_data in enumerate(data_array):
                        if row_idx >= rows:
                            break
                        for col_idx, cell_data in enumerate(row_data):
                            if col_idx >= cols:
                                break
                            table.cell(row_idx, col_idx).text = _process_text_escapes(str(cell_data))
                except json.JSONDecodeError:
                    return "Error: Invalid JSON format for data parameter"

            state.is_modified = True
            return f"Successfully added {rows}x{cols} table on slide {slide_number}\nShape ID: {shape.shape_id}"
        except Exception as e:
            return f"Error adding table: {str(e)}"

    @mcp.tool()
    def add_chart(
        slide_number: int,
        chart_type: str,
        categories: str,
        series_data: str,
        left: float = 1.0,
        top: float = 2.0,
        width: float = 8.0,
        height: float = 4.5
    ) -> str:
        """
        Add a chart to a slide.

        Args:
            slide_number: Target slide number (1-based)
            chart_type: Chart type - "bar", "column", "line", "pie", "area"
            categories: JSON array of category labels, e.g., '["Q1", "Q2", "Q3", "Q4"]'
            series_data: JSON object with series, e.g.,
                        '{"Sales": [100, 120, 140, 160], "Profit": [20, 25, 30, 35]}'
                        For pie charts, use single series: '{"Market Share": [40, 30, 20, 10]}'
            left: Left position in inches (default 1.0)
            top: Top position in inches (default 2.0)
            width: Width in inches (default 8.0)
            height: Height in inches (default 4.5)

        Returns:
            Success message with shape ID, or error message
        """
        if state.presentation is None:
            return "Error: No presentation is currently open"

        prs = state.presentation
        if slide_number < 1 or slide_number > len(prs.slides):
            return f"Error: slide_number {slide_number} is out of range (1-{len(prs.slides)})"

        slide = prs.slides[slide_number - 1]

        # Map chart types
        chart_map = {
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
            "area": XL_CHART_TYPE.AREA,
        }

        chart_type_lower = chart_type.lower()
        if chart_type_lower not in chart_map:
            return f"Error: Unknown chart type '{chart_type}'. Valid types: {', '.join(chart_map.keys())}"

        try:
            # Parse data
            cats = json.loads(categories)
            series = json.loads(series_data)

            # Build chart data
            chart_data = CategoryChartData()
            chart_data.categories = cats

            for series_name, values in series.items():
                chart_data.add_series(series_name, values)

            # Add chart
            x, y, cx, cy = Inches(left), Inches(top), Inches(width), Inches(height)
            chart = slide.shapes.add_chart(
                chart_map[chart_type_lower], x, y, cx, cy, chart_data
            ).chart

            state.is_modified = True
            return f"Successfully added {chart_type} chart on slide {slide_number}"
        except json.JSONDecodeError as e:
            return f"Error: Invalid JSON format - {str(e)}"
        except Exception as e:
            return f"Error adding chart: {str(e)}"


def _parse_color(hex_color: str) -> RGBColor:
    """Parse a hex color string to RGBColor.

    Args:
        hex_color: Color as hex string (e.g., "#FF0000" or "FF0000")

    Returns:
        RGBColor object, or None if parsing fails
    """
    try:
        hex_color = hex_color.lstrip('#')
        if len(hex_color) == 6:
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return RGBColor(r, g, b)
        else:
            print(f"Warning: Invalid hex color '{hex_color}' - expected 6 characters")
    except ValueError as e:
        print(f"Warning: Could not parse hex color '{hex_color}': {e}")
    return None
